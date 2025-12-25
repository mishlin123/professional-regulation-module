import json
from pptx import Presentation

def extract_content(pptx_path, output_path):
    prs = Presentation(pptx_path)
    slides_data = []

    for i, slide in enumerate(prs.slides):
        slide_content = {
            "id": i + 1,
            "title": "",
            "body_text": [],
            "notes": ""
        }

        # Extract Title
        if slide.shapes.title:
            slide_content["title"] = slide.shapes.title.text

        # Extract Text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape != slide.shapes.title:
                if shape.text.strip():
                    slide_content["body_text"].append(shape.text.strip())

        # Extract Notes
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes:
                slide_content["notes"] = notes.strip()

        slides_data.append(slide_content)

    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(slides_data, f, indent=2, ensure_ascii=False)

    print(f"Extracted {len(slides_data)} slides to {output_path}")

if __name__ == "__main__":
    extract_content("resources/course.pptx", "resources/course_content.json")
